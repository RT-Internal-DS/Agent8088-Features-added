"""Text extraction from office/PDF documents, so the model can read them as text.

Called from the file-read path as a pre-step: if `extract_text` recognizes the
extension it returns extracted text (or an error string), otherwise it returns
None and the caller falls back to reading the file as plain text/bytes.

.docx/.pptx need no dependency — both are zip files of XML, and the parts we
care about (`word/document.xml`, `ppt/slides/slideN.xml`) are handled with
stdlib `zipfile` + `xml.etree.ElementTree`. .xlsx and .pdf go through
`openpyxl`/`pypdf`, which are optional dependencies imported lazily so this
module (and anything that imports it) loads fine without them.
"""
from __future__ import annotations

import re
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

# Caller's normal file-read path caps at ~2MB; that cap doesn't apply here, so
# without an explicit guard a crafted large document is an unbounded-memory
# read. 25MB is generous for a document someone actually wants summarized.
MAX_DOCUMENT_BYTES = 25 * 1024 * 1024

# Independent of the input-size guard above: a small zip can still decompress
# to gigabytes of XML (zip bomb / entity expansion). Extraction stops once the
# accumulated output text passes this, regardless of how much of the document
# is left unread.
MAX_EXTRACTED_CHARS = 5 * 1024 * 1024

WORD_NS = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
DRAWING_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# ET.parse uses Python's default (non-network) XML parser: it has no DTD
# support, so it cannot fetch or resolve external entities. That closes the
# classic XXE hole. It does not defend against a purely-internal entity-expansion
# bomb (billion laughs) — nothing in the stdlib does without a third-party
# parser — so MAX_EXTRACTED_CHARS above is the actual backstop against that.
# These parts always come from a local zip we just opened, never from network
# input, which narrows the exposure further.


def extract_text(path, max_bytes: int = MAX_DOCUMENT_BYTES):
    """Extract readable text from a document. Return None if `path` is not a
    document format this module handles, so the caller falls through to normal
    text reading.
    """
    path = Path(path)
    ext = path.suffix.lower()
    if ext not in (".docx", ".pptx", ".xlsx", ".pdf"):
        return None

    size = path.stat().st_size
    if size > max_bytes:
        raise ValueError(
            f"Document is too large to read (limit: {max_bytes} bytes): {path}"
        )

    if ext == ".docx":
        return _extract_docx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".xlsx":
        return _extract_xlsx(path)
    return _extract_pdf(path)


def _truncated(chunks) -> str:
    """Join chunks, stopping once past MAX_EXTRACTED_CHARS rather than joining
    everything first — a zip bomb should stop being decompressed, not just be
    truncated after the fact."""
    out = []
    total = 0
    for chunk in chunks:
        out.append(chunk)
        total += len(chunk)
        if total > MAX_EXTRACTED_CHARS:
            out.append("\n[... truncated: extracted text exceeded limit ...]")
            break
    return "\n".join(out)


# ---------------------------------------------------------------------------
# .docx
# ---------------------------------------------------------------------------
def _extract_docx(path):
    try:
        with zipfile.ZipFile(path) as zf:
            with zf.open("word/document.xml") as f:
                root = ET.parse(f).getroot()
    except zipfile.BadZipFile:
        return f"Could not read {path}: not a valid .docx (bad zip)."
    except KeyError:
        return f"Could not read {path}: missing word/document.xml (not a valid .docx)."

    body = root.find(f"{WORD_NS}body")
    if body is None:
        return ""

    def chunks():
        for el in body:
            tag = el.tag
            if tag == f"{WORD_NS}p":
                text = "".join(t.text or "" for t in el.iter(f"{WORD_NS}t"))
                yield text
            elif tag == f"{WORD_NS}tbl":
                for tr in el.iter(f"{WORD_NS}tr"):
                    cells = []
                    for tc in tr.findall(f"{WORD_NS}tc"):
                        cells.append("".join(t.text or "" for t in tc.iter(f"{WORD_NS}t")))
                    yield "| " + " | ".join(cells) + " |"

    return _truncated(chunks())


# ---------------------------------------------------------------------------
# .pptx
# ---------------------------------------------------------------------------
def _slide_number(name):
    # "slide2.xml" -> 2, "slide10.xml" -> 10 — naive string sort would put
    # slide10 before slide2.
    match = re.search(r"(\d+)", name)
    return int(match.group(1)) if match else 0


def _extract_pptx(path):
    try:
        with zipfile.ZipFile(path) as zf:
            names = zf.namelist()
            slides = sorted(
                (n for n in names if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)),
                key=_slide_number,
            )
            if not slides:
                return f"Could not read {path}: no slides found (not a valid .pptx)."

            def chunks():
                for slide_name in slides:
                    n = _slide_number(slide_name)
                    yield f"## Slide {n}"
                    with zf.open(slide_name) as f:
                        root = ET.parse(f).getroot()
                    yield "\n".join(t.text or "" for t in root.iter(f"{DRAWING_NS}t"))

                    notes_name = f"ppt/notesSlides/notesSlide{n}.xml"
                    if notes_name in names:
                        with zf.open(notes_name) as f:
                            notes_root = ET.parse(f).getroot()
                        notes_text = "\n".join(
                            t.text or "" for t in notes_root.iter(f"{DRAWING_NS}t")
                        ).strip()
                        if notes_text:
                            yield f"Notes: {notes_text}"

            return _truncated(chunks())
    except zipfile.BadZipFile:
        return f"Could not read {path}: not a valid .pptx (bad zip)."
    except KeyError:
        return f"Could not read {path}: malformed .pptx (missing expected part)."


# ---------------------------------------------------------------------------
# .xlsx
# ---------------------------------------------------------------------------
def _extract_xlsx(path):
    try:
        import openpyxl
    except ImportError:
        return "openpyxl is not installed. Install it with:\n  pip install openpyxl"

    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    except Exception as exc:
        return f"Could not read {path}: {exc}"

    def chunks():
        for sheet in wb.worksheets:
            yield f"## Sheet: {sheet.title}"
            for row in sheet.iter_rows():
                # data_only=True still leaves None for a formula cell with no
                # cached value — render that as empty, not the string "None".
                cells = ["" if c.value is None else str(c.value) for c in row]
                if not any(cells):
                    continue  # skip fully-empty rows (also drops empty trailing ones)
                yield "| " + " | ".join(cells) + " |"

    try:
        return _truncated(chunks())
    finally:
        wb.close()


# ---------------------------------------------------------------------------
# .pdf
# ---------------------------------------------------------------------------
def _extract_pdf(path):
    try:
        import pypdf
    except ImportError:
        return "pypdf is not installed. Install it with:\n  pip install pypdf"

    try:
        reader = pypdf.PdfReader(path)
    except Exception as exc:
        return f"Could not read {path}: {exc}"

    extracted = []

    def chunks():
        for i, page in enumerate(reader.pages, start=1):
            yield f"## Page {i}"
            try:
                text = page.extract_text() or ""
            except Exception as exc:
                yield f"[could not extract page {i}: {exc}]"
                continue
            extracted.append(text)
            yield text

    out = _truncated(chunks())
    # A scanned PDF is a stack of images: pypdf finds no text layer and returns
    # empty strings for every page. Saying so is the honest answer — silently
    # returning a page-header skeleton reads as "this document is blank", which
    # is the one conclusion that is definitely wrong.
    if not any(t.strip() for t in extracted):
        return (f"{path} has {len(reader.pages)} page(s) but no extractable text. "
                "It is most likely a scan or image-only PDF; reading it would need OCR.")
    return out


# ---------------------------------------------------------------------------
# Writing
# ---------------------------------------------------------------------------
# The line format is deliberately flat rather than nested JSON: this exists as
# the fallback for when generating openpyxl/python-docx code defeats a smaller
# model, and a model that cannot reliably emit balanced JSON can still emit
# "# Heading" one line at a time. Anything richer belongs in the `documents`
# skill, which teaches writing the library code directly.
BUILDABLE = (".docx", ".xlsx", ".pptx")


def build_document(path, content: str) -> str:
    """Create a .docx/.xlsx/.pptx from a flat line format. Returns a summary.

    docx:  "# Heading", "## Subheading", "| a | b |" table rows, else paragraph
    xlsx:  "## Sheet: Name" starts a sheet, "| a | b |" rows are cells
    pptx:  "## Slide title" starts a slide, other lines are bullets
    """
    path = Path(path)
    ext = path.suffix.lower()
    lines = (content or "").splitlines()
    if ext == ".docx":
        return _build_docx(path, lines)
    if ext == ".xlsx":
        return _build_xlsx(path, lines)
    if ext == ".pptx":
        return _build_pptx(path, lines)
    return (f"create_document cannot build '{ext or path.name}'. "
            f"Supported: {', '.join(BUILDABLE)}. For PDF or anything needing "
            f"layout control, write the file with reportlab through execute_shell.")


def _cells(line):
    return [c.strip() for c in line.strip().strip("|").split("|")]


def _build_docx(path, lines):
    try:
        import docx
    except ImportError:
        return "python-docx is not installed. Install it with:\n  pip install python-docx"
    doc = docx.Document()
    counts = {"headings": 0, "paragraphs": 0, "tables": 0}
    rows = []

    def flush():
        # Consecutive "| a | b |" lines are one table, so buffer them and emit
        # on the first non-row line instead of making a table per row.
        if not rows:
            return
        table = doc.add_table(rows=len(rows), cols=max(len(r) for r in rows))
        table.style = "Table Grid"
        for i, row in enumerate(rows):
            for j, cell in enumerate(row):
                table.rows[i].cells[j].text = cell
        counts["tables"] += 1
        rows.clear()

    for line in lines:
        stripped = line.strip()
        if stripped.startswith("|") and stripped.endswith("|"):
            rows.append(_cells(stripped))
            continue
        flush()
        if stripped.startswith("## "):
            doc.add_heading(stripped[3:], level=2)
            counts["headings"] += 1
        elif stripped.startswith("# "):
            doc.add_heading(stripped[2:], level=1)
            counts["headings"] += 1
        elif stripped:
            doc.add_paragraph(stripped)
            counts["paragraphs"] += 1
    flush()
    doc.save(str(path))
    return (f"Created {path} — {counts['headings']} heading(s), "
            f"{counts['paragraphs']} paragraph(s), {counts['tables']} table(s)")


def _build_xlsx(path, lines):
    try:
        import openpyxl
    except ImportError:
        return "openpyxl is not installed. Install it with:\n  pip install openpyxl"
    wb = openpyxl.Workbook()
    sheet = wb.active
    first = True
    sheets = rows = 0
    for line in lines:
        stripped = line.strip()
        if stripped.lower().startswith("## sheet:"):
            title = stripped.split(":", 1)[1].strip() or "Sheet"
            if first:
                sheet.title, first = title, False
            else:
                sheet = wb.create_sheet(title)
            sheets += 1
        elif stripped.startswith("|") and stripped.endswith("|"):
            # A number written as text sorts and charts wrong in Excel, so
            # values that look numeric are stored as numbers.
            out = []
            for cell in _cells(stripped):
                try:
                    out.append(float(cell) if "." in cell else int(cell))
                except ValueError:
                    out.append(cell)
            sheet.append(out)
            rows += 1
    wb.save(str(path))
    return f"Created {path} — {max(sheets, 1)} sheet(s), {rows} row(s)"


def _build_pptx(path, lines):
    try:
        from pptx import Presentation
    except ImportError:
        return "python-pptx is not installed. Install it with:\n  pip install python-pptx"
    prs = Presentation()
    slide = None
    slides = bullets = 0
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("## "):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = stripped[3:]
            slides += 1
        elif stripped and slide is not None:
            body = slide.placeholders[1].text_frame
            para = body.paragraphs[0] if not body.text else body.add_paragraph()
            para.text = stripped
            bullets += 1
    if slides == 0:
        return ("No slides created: a .pptx needs at least one '## Slide title' "
                "line to start a slide.")
    prs.save(str(path))
    return f"Created {path} — {slides} slide(s), {bullets} bullet(s)"


# ---------------------------------------------------------------------------
# Conversion — LibreOffice headless
# ---------------------------------------------------------------------------
# Skill-only guidance ("run soffice via execute_shell") proved unreliable
# against the weak local model this project targets: given "convert X to
# PDF", it wrote a fresh reportlab script generating an unrelated document
# from scratch rather than following the documented soffice command, twice,
# even after the skill was rewritten to lead with that instruction. A
# deterministic tool the model just calls with two arguments removes the
# chance to substitute its own approach — same reasoning as build_document.
CONVERTIBLE_TARGETS = ("pdf", "docx", "pptx", "xlsx")

# LibreOffice headless opens each input format as a specific app, and can only
# export to that app's formats (plus PDF, which every app exports). PDF is the
# exception: it opens as a Draw document, and Draw does NOT export to
# Writer/Impress/Calc formats — so "convert PDF to docx" fails with
# "no export filter" every time. Mapping this keeps the tool from claiming a
# conversion is possible, running soffice, and returning a raw filter error
# that the model then tries to work around with a cascade of shell commands.
_WRITER_FORMATS = (".docx", ".doc", ".odt", ".rtf")
_CALC_FORMATS = (".xlsx", ".xls", ".csv", ".ods")
_IMPRESS_FORMATS = (".pptx", ".ppt", ".odp")
_DRAW_FORMATS = (".pdf",)  # PDF imports as a Draw doc; Draw exports to pdf/html/png, not docx/pptx/xlsx

_SOFFICE_INSTALL_PATHS = (
    r"C:\Program Files\LibreOffice\program\soffice.exe",
    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
)


def _soffice_executable():
    """Find soffice, or None. Checked fresh every call — the installer's
    LibreOffice step is best-effort and can have failed or been skipped."""
    import shutil
    found = shutil.which("soffice")
    if found:
        return found
    for candidate in _SOFFICE_INSTALL_PATHS:
        if Path(candidate).exists():
            return candidate
    return None


def _conversion_is_supported(src_ext: str, target_format: str) -> bool:
    """Whether LibreOffice headless can convert this source to this target.

    Every app can export to PDF. Within an app family, the native formats
    convert to each other. PDF is the narrow case: it imports as Draw, and
    Draw only exports to PDF (and image/html formats we don't expose) — not
    to docx/pptx/xlsx, which belong to other apps.
    """
    src_ext = src_ext.lower()
    target_format = target_format.lower()
    if target_format == "pdf":
        return True  # every app exports to PDF
    if src_ext in _DRAW_FORMATS:
        return False  # PDF source can only go to PDF, not to editable formats
    if src_ext in _WRITER_FORMATS:
        return target_format in ("docx",)  # expose only the modern Writer target
    if src_ext in _CALC_FORMATS:
        return target_format in ("xlsx",)
    if src_ext in _IMPRESS_FORMATS:
        return target_format in ("pptx",)
    return False  # unknown source family


def convert_document(path, target_format: str, timeout: int = 60) -> str:
    """Convert `path` to `target_format` via LibreOffice headless, in place
    (same directory, same basename, new extension). Returns a summary or a
    plain-language reason it didn't happen — never raises, so a tool caller
    can return this string directly as the result.
    """
    path = Path(path)
    target_format = (target_format or "").strip().lower().lstrip(".")
    if target_format not in CONVERTIBLE_TARGETS:
        return (f"Cannot convert to '{target_format}'. Supported targets: "
                 f"{', '.join(CONVERTIBLE_TARGETS)}.")

    src_ext = path.suffix.lower()
    if not _conversion_is_supported(src_ext, target_format):
        # Be explicit: name the unsupported pair and what IS supported for this
        # source, so the model does not fall back to a shell workaround cascade.
        if src_ext == ".pdf":
            hint = ("PDF can only be re-exported as PDF via LibreOffice. "
                    "To extract PDF text into an editable document, use a "
                    "PDF-to-text tool or library (e.g. pdfplumber) instead.")
        else:
            hint = f"Convert {src_ext} files to one of: {', '.join(CONVERTIBLE_TARGETS)}."
        return (f"Conversion not supported: {src_ext} -> .{target_format}. {hint}")

    soffice = _soffice_executable()
    if not soffice:
        return ("LibreOffice is not installed, so conversion is unavailable. "
                 "Install it with: winget install TheDocumentFoundation.LibreOffice "
                 "(or rerun the Agent8088 installer), then try again.")

    if not path.exists():
        return f"Cannot convert: {path} does not exist."

    output_path = path.with_suffix("." + target_format)
    import subprocess
    try:
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", target_format,
             "--outdir", str(path.parent), str(path)],
            capture_output=True, text=True, timeout=timeout,
        )
    except subprocess.TimeoutExpired:
        return (f"LibreOffice timed out after {timeout}s converting {path.name}. "
                 "A first-run profile setup can be slow — try again.")

    # soffice can print a success-looking line and still not produce the file
    # (locked output, unsupported filter for this input) — check disk, not
    # the exit code or stdout text, before claiming success.
    if not output_path.exists():
        detail = (result.stderr or result.stdout or "no output from soffice").strip()
        return (f"Conversion failed: {path.name} was not converted to "
                 f"{target_format}. LibreOffice said: {detail[:300]}")

    return f"Converted {path.name} to {output_path.name} ({output_path.stat().st_size} bytes)."


if __name__ == "__main__":
    # Minimal self-check: build a real .docx/.pptx in a temp dir and confirm
    # extraction round-trips. Run with: python -m agent8088.documents
    import tempfile

    with tempfile.TemporaryDirectory() as tmp:
        tmp = Path(tmp)

        # .docx: one paragraph + one 2x2 table
        docx_path = tmp / "t.docx"
        doc_xml = (
            '<?xml version="1.0"?>'
            '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            "<w:body>"
            '<w:p><w:r><w:t>Hello world</w:t></w:r></w:p>'
            "<w:tbl>"
            '<w:tr><w:tc><w:p><w:r><w:t>a</w:t></w:r></w:p></w:tc>'
            '<w:tc><w:p><w:r><w:t>b</w:t></w:r></w:p></w:tc></w:tr>'
            "</w:tbl>"
            "</w:body></w:document>"
        )
        with zipfile.ZipFile(docx_path, "w") as zf:
            zf.writestr("word/document.xml", doc_xml)
        out = extract_text(docx_path)
        assert "Hello world" in out and "| a | b |" in out, out

        # .pptx: slide2 and slide10 must sort numerically, not lexically
        pptx_path = tmp / "t.pptx"
        with zipfile.ZipFile(pptx_path, "w") as zf:
            for n, text in ((2, "second"), (10, "tenth")):
                slide_xml = (
                    '<?xml version="1.0"?>'
                    '<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
                    'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                    f"<a:t>{text}</a:t></p:sld>"
                )
                zf.writestr(f"ppt/slides/slide{n}.xml", slide_xml)
        out = extract_text(pptx_path)
        assert out.index("## Slide 2") < out.index("## Slide 10"), out
        assert out.index("second") < out.index("tenth"), out

        # unknown extension -> None
        txt_path = tmp / "t.txt"
        txt_path.write_text("plain")
        assert extract_text(txt_path) is None

        # corrupt docx -> message, not a crash
        bad_path = tmp / "bad.docx"
        bad_path.write_text("not a zip")
        result = extract_text(bad_path)
        assert result is not None and "not a valid .docx" in result

        # size guard
        big_path = tmp / "big.docx"
        big_path.write_bytes(b"0" * 100)
        try:
            extract_text(big_path, max_bytes=10)
            raise AssertionError("expected ValueError")
        except ValueError:
            pass

    print("documents.py self-check passed")
